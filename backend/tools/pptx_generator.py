"""Branded PowerPoint pitch deck generator.

Generates a professional .pptx from slide sections + company branding.
Used by the Funding Kit and available as a tool for agents.
"""
from __future__ import annotations

import io
import logging
import os
import tempfile
import uuid
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

_DEFAULT_PRIMARY = "#002EFF"
_DEFAULT_TEXT = "#111827"
_SLIDE_W = 9144000   # 10 inches in EMU (914400 per inch)
_SLIDE_H = 6858000   # 7.5 inches


def _hex_rgb(hex_color: str) -> tuple[int, int, int]:
    h = (hex_color or _DEFAULT_PRIMARY).lstrip("#")
    if len(h) != 6:
        return (0, 46, 255)
    try:
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except ValueError:
        return (0, 46, 255)


def _is_dark(r: int, g: int, b: int) -> bool:
    return (0.299 * r + 0.587 * g + 0.114 * b) < 128


def _fetch_logo_bytes(logo_url_or_path: str) -> bytes | None:
    if not logo_url_or_path:
        return None
    try:
        p = Path(logo_url_or_path)
        if p.exists() and p.is_file():
            return p.read_bytes()
        if logo_url_or_path.startswith("http"):
            import requests
            resp = requests.get(logo_url_or_path, timeout=5)
            if resp.status_code == 200:
                return resp.content
    except Exception as exc:
        logger.warning("pptx: could not fetch logo %s: %s", logo_url_or_path, exc)
    return None


def _safe_text(text: str) -> str:
    _MAP = {
        "—": "--", "–": "-", "‘": "'", "’": "'",
        "“": '"', "”": '"', "…": "...", "•": "*",
        " ": " ",
    }
    for ch, rep in _MAP.items():
        text = text.replace(ch, rep)
    return text


def generate_pptx(
    title: str = "",
    slides: list[dict] | None = None,
    sections: list[dict] | None = None,
    company_name: str = "",
    primary_color: str = "",
    accent_color: str = "",
    logo_url: str = "",
    logo_path: str = "",
    filename: str = "",
    output_dir: str = "",
    founder_id: str = "",
    **kwargs,
) -> dict[str, Any]:
    """Generate a branded PowerPoint .pptx.

    Args:
        title: Presentation title (cover slide heading).
        slides / sections: list of {heading, body} dicts — one per slide.
        company_name: Shown in footer of every slide.
        primary_color: Hex color for cover + accent elements (#RRGGBB).
        logo_url / logo_path: URL or path to company logo image.
        filename: Output filename (will be given a unique suffix).
        founder_id: If provided, auto-load branding from company genome.
    Returns:
        {generated: bool, path: str, filename: str} or {generated: False, error: str}
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.oxml.ns import qn
        from lxml import etree
    except ImportError as e:
        return {"generated": False, "error": f"python-pptx not installed: {e}"}

    # ── Resolve branding from genome if founder_id supplied ───────────────────
    if founder_id and not primary_color:
        try:
            from backend.genome.store import get_genome
            genome = get_genome(founder_id) or {}
            branding = _extract_branding(genome)
            primary_color = primary_color or branding.get("primary_color", "")
            accent_color = accent_color or branding.get("accent_color", "")
            logo_url = logo_url or branding.get("logo_url", "")
            logo_path = logo_path or branding.get("logo_path", "")
            company_name = company_name or branding.get("company_name", "")
        except Exception as exc:
            logger.warning("pptx: genome branding lookup failed: %s", exc)

    primary_color = primary_color or _DEFAULT_PRIMARY
    pr, pg, pb = _hex_rgb(primary_color)
    cover_text_color = (255, 255, 255) if _is_dark(pr, pg, pb) else (17, 24, 39)

    # Light tint for content slide backgrounds (very light version of primary)
    accent_light = (
        min(255, pr + int((255 - pr) * 0.92)),
        min(255, pg + int((255 - pg) * 0.92)),
        min(255, pb + int((255 - pb) * 0.92)),
    )

    # ── Resolve output path ───────────────────────────────────────────────────
    if not output_dir:
        vault = os.environ.get("OBSIDIAN_VAULT", "/tmp/astra_docs")
        output_dir = str(vault).rstrip("/") + "/files"
    os.makedirs(output_dir, exist_ok=True)

    if filename:
        safe_name = Path(filename).name.encode("ascii", "ignore").decode()
        if not safe_name.lower().endswith(".pptx"):
            safe_name = safe_name.rsplit(".", 1)[0] + ".pptx"
        stem = safe_name.rsplit(".", 1)[0]
        filename = f"{stem}_{uuid.uuid4().hex[:6]}.pptx"
    else:
        safe_title = (title or "presentation").lower().replace(" ", "_").encode("ascii", "ignore").decode()
        filename = f"{safe_title}_{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join(output_dir, filename)

    slide_list = slides or sections or []

    # ── Logo bytes ─────────────────────────────────────────────────────────────
    logo_source = logo_path or logo_url
    if not logo_source and founder_id:
        # Also check vault for a logo file
        vault = os.environ.get("OBSIDIAN_VAULT", "/tmp/astra_docs")
        for ext in ("png", "jpg", "jpeg", "webp"):
            candidate = Path(vault) / f"logo.{ext}"
            if candidate.exists():
                logo_source = str(candidate)
                break
    logo_bytes = _fetch_logo_bytes(logo_source)

    # ── Build presentation ─────────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width = Emu(_SLIDE_W)
    prs.slide_height = Emu(_SLIDE_H)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ── Cover slide ────────────────────────────────────────────────────────────
    cover = prs.slides.add_slide(blank_layout)

    # Full-bleed background rectangle
    bg = cover.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(0), Emu(0), Emu(_SLIDE_W), Emu(_SLIDE_H),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(*[pr, pg, pb])
    bg.line.fill.background()

    # Logo top-left if available
    _add_logo(cover, logo_bytes, Inches(0.4), Inches(0.3), max_h=Inches(0.7))

    # Company name (small, top line)
    if company_name:
        _add_textbox(
            cover, _safe_text(company_name),
            left=Inches(0.4), top=Inches(1.1), width=Inches(9.2), height=Inches(0.5),
            font_size=Pt(13), bold=False,
            color=RGBColor(*(int(c * 0.75) for c in cover_text_color)),
            align=PP_ALIGN.LEFT,
        )

    # Main title
    _add_textbox(
        cover, _safe_text(title or company_name or "Pitch Deck"),
        left=Inches(0.6), top=Inches(2.0), width=Inches(8.8), height=Inches(1.8),
        font_size=Pt(44), bold=True,
        color=RGBColor(*cover_text_color),
        align=PP_ALIGN.LEFT,
    )

    # Subtitle / label
    _add_textbox(
        cover, "Investor Presentation",
        left=Inches(0.6), top=Inches(3.9), width=Inches(8.0), height=Inches(0.5),
        font_size=Pt(18), bold=False,
        color=RGBColor(*(min(255, int(c * 0.85)) for c in cover_text_color)),
        align=PP_ALIGN.LEFT,
    )

    # Accent bottom bar
    bar = cover.shapes.add_shape(1, Emu(0), Emu(_SLIDE_H - 457200), Emu(_SLIDE_W), Emu(457200))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(
        max(0, pr - 30), max(0, pg - 30), max(0, pb - 30)
    )
    bar.line.fill.background()

    # ── Content slides ─────────────────────────────────────────────────────────
    for i, section in enumerate(slide_list):
        if isinstance(section, str):
            section = {"heading": "", "body": section}
        heading = _safe_text(str(section.get("heading") or section.get("title") or ""))
        body = _safe_text(str(section.get("body") or section.get("content") or ""))
        if not heading and not body:
            continue

        sl = prs.slides.add_slide(blank_layout)

        # White background (already default but be explicit)
        bg2 = sl.shapes.add_shape(1, Emu(0), Emu(0), Emu(_SLIDE_W), Emu(_SLIDE_H))
        bg2.fill.solid()
        bg2.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg2.line.fill.background()

        # Left accent bar
        bar2 = sl.shapes.add_shape(1, Emu(0), Emu(0), Emu(182880), Emu(_SLIDE_H))
        bar2.fill.solid()
        bar2.fill.fore_color.rgb = RGBColor(pr, pg, pb)
        bar2.line.fill.background()

        # Top header band (light tint)
        header = sl.shapes.add_shape(1, Emu(182880), Emu(0), Emu(_SLIDE_W - 182880), Emu(950000))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(*accent_light)
        header.line.fill.background()

        # Logo in header top-right
        _add_logo(sl, logo_bytes, Inches(8.5), Inches(0.12), max_h=Inches(0.6))

        # Slide heading
        _add_textbox(
            sl, heading,
            left=Inches(0.45), top=Inches(0.12), width=Inches(8.0), height=Inches(0.75),
            font_size=Pt(22), bold=True,
            color=RGBColor(pr, pg, pb),
            align=PP_ALIGN.LEFT,
        )

        # Body text
        _add_textbox(
            sl, body,
            left=Inches(0.45), top=Inches(1.1), width=Inches(9.0), height=Inches(5.5),
            font_size=Pt(14), bold=False,
            color=RGBColor(17, 24, 39),
            align=PP_ALIGN.LEFT,
        )

        # Footer: company name + slide number
        footer_text = f"{company_name}  |  {i + 1}" if company_name else str(i + 1)
        _add_textbox(
            sl, footer_text,
            left=Inches(0.45), top=Inches(6.9), width=Inches(9.0), height=Inches(0.35),
            font_size=Pt(9), bold=False,
            color=RGBColor(150, 150, 150),
            align=PP_ALIGN.LEFT,
        )

    prs.save(filepath)
    return {"generated": True, "path": filepath, "filename": Path(filepath).name}


# ── Helpers ────────────────────────────────────────────────────────────────────

def _add_textbox(slide, text, left, top, width, height, font_size, bold, color, align):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _add_logo(slide, logo_bytes: bytes | None, left, top, max_h):
    if not logo_bytes:
        return
    try:
        from pptx.util import Emu
        import io
        img_stream = io.BytesIO(logo_bytes)
        pic = slide.shapes.add_picture(img_stream, left, top)
        # Scale to max_h keeping aspect ratio
        if pic.height > max_h:
            ratio = max_h / pic.height
            pic.width = int(pic.width * ratio)
            pic.height = int(max_h)
    except Exception as exc:
        logger.warning("pptx: logo embed failed: %s", exc)


def _extract_branding(genome: dict) -> dict:
    """Pull branding fields out of a genome dict."""
    sections = genome.get("sections") or {}
    branding = sections.get("branding") or {}
    profile = sections.get("profile") or {}
    brand_voice = sections.get("brand_voice") or {}

    def _val(d: dict, key: str) -> str:
        entry = d.get(key)
        if isinstance(entry, dict):
            return str(entry.get("value") or "")
        return str(entry or "")

    # Company name
    name = (
        _val(profile, "name")
        or _val(profile, "company_name")
        or genome.get("company_name")
        or ""
    )
    # Primary color: check branding section first, then brand_voice
    primary = (
        _val(branding, "primary_color")
        or _val(branding, "color")
        or _val(brand_voice, "primary_color")
        or _val(brand_voice, "color")
        or ""
    )
    accent = _val(branding, "accent_color") or _val(brand_voice, "accent_color") or ""
    logo_url = _val(branding, "logo_url") or _val(branding, "logo") or ""
    logo_path = _val(branding, "logo_path") or ""

    return {
        "company_name": name,
        "primary_color": primary,
        "accent_color": accent,
        "logo_url": logo_url,
        "logo_path": logo_path,
    }
